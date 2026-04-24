import io
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader
from langchain_text_splitters import RecursiveCharacterTextSplitter


class DocumentParser:
    def __init__(self):
        pass

    def parse(self, file_bytes, mime_type):
        if mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
            return self._parse_docx(file_bytes)
        elif mime_type == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet':
            return self._parse_xlsx(file_bytes)
        elif mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
            return self._parse_pptx(file_bytes)
        elif mime_type == 'application/pdf':
            return self._parse_pdf(file_bytes)
        elif mime_type.startswith('image/'):
            return self._parse_image(file_bytes)
        else:
            raise ValueError(f"Unsupported MIME type: {mime_type}")

    def _parse_docx(self, file_bytes):
        doc = Document(io.BytesIO(file_bytes))
        return "\n".join([para.text for para in doc.paragraphs])

    def _parse_xlsx(self, file_bytes):
        workbook = load_workbook(io.BytesIO(file_bytes), data_only=True)
        data = []
        for sheet in workbook:
            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                if row_text:
                    data.append(row_text)
        return "\n".join(data)

    def _parse_pptx(self, file_bytes):
        prs = Presentation(io.BytesIO(file_bytes))
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    def _parse_pdf(self, file_bytes):
        reader = PdfReader(io.BytesIO(file_bytes))
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n"
        return text

    def _parse_image(self, file_bytes):
        return ""


class TextChunker:
    def __init__(self, config):
        self.chunk_size = config.chunk_size
        self.chunk_overlap = config.chunk_overlap
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=self.chunk_size,
            chunk_overlap=self.chunk_overlap
        )

    def chunk_text(self, text, metadata=None):
        return self.text_splitter.create_documents(
            [text],
            metadatas=[metadata] if metadata else None
        )


class IngestionPipeline:
    def __init__(self, config, drive_client, registry, parser, chunker, vector_store):
        self.config = config
        self.drive_client = drive_client
        self.registry = registry
        self.parser = parser
        self.chunker = chunker
        self.vector_store = vector_store

    def run_delta_check(self):
        drive_files = self.drive_client.list_files_in_folder(self.config.drive_folder_id)
        new_or_updated_files = self.registry.get_new_or_updated_files(drive_files)
        return new_or_updated_files

    def ingest_file(self, file):
        file_id = file['id']
        modified_time = file['modifiedTime']
        metadata = {
            'file_id': file_id,
            'file_name': file['name'],
            'mime_type': file['mimeType'],
            'modified_time': modified_time
        }
        file_bytes = self.drive_client.download_file(file_id).read()
        text = self.parser.parse(file_bytes, file['mimeType'])
        chunks = self.chunker.chunk_text(text, metadata)
        self.vector_store.upsert(chunks)
        self.registry.update(file_id, modified_time)

    def run(self):
        new_or_updated_files = self.run_delta_check()
        for file in new_or_updated_files:
            print(f"Ingesting: {file['name']}")
            self.ingest_file(file)
        print(f"Ingestion complete. {len(new_or_updated_files)} files processed.")
